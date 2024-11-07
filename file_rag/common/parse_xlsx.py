from pptx import Presentation
from pptx.shapes.picture import Picture

prs = Presentation("E:\\vs-python\\AGENTS\\data\\1930ff7e248a2000\\main_doc\\政务OA大模型.pptx")
index = 1
#读取幻灯片的每一页
for slide in prs.slides:
    # 读取每一板块
    for shape in slide.shapes:
        # print(dir(shape))
        #是否有文字框
        if shape.has_text_frame:
            #读文字框的每一段落
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text:
                    # 输出段落文字,也有一些属性,可以用dir查看
                    # print(dir(paragraph))
                    print(paragraph.text)
        #是否有表格
        elif shape.has_table:
            one_table_data = []
            for row in shape.table.rows:  # 读每行
                row_data = []
                for cell in row.cells:  # 读一行中的所有单元格
                    c = cell.text
                    row_data.append(c)
                one_table_data.append(row_data)  # 把每一行存入表
            #用二维列表输出表格行和列的数据
            print(one_table_data)
        # 是否有图片
        elif isinstance(shape, Picture):
            #shape.image.blob:二进制图像字节流,写入图像文件
            with open(f'{index}.jpg', 'wb') as f:
                f.write(shape.image.blob)
                index += 1

