import os
from PIL import Image
import re
import sys
from pathlib import Path
sys.path.append('./')
import pandas as pd
from PyPDF2 import PdfReader
from docx import Document
import pytesseract
import xlrd
import fitz
import io
from pprint import pprint
from pptx import Presentation
from pptx.shapes.picture import Picture
from file_rag.settings import Settings
from file_rag.knowledge_base.utils import get_file_path, get_kb_path
from file_rag.server.utils import api_address
from file_rag.webui_pages.utils import ApiRequest

api_base_url = api_address()
api: ApiRequest = ApiRequest(api_base_url)
# 定义文本提取函数
    # '''
    #  该方法支持 pdf、pdf扫描件、docx、doc、wps、xlsx、pptx等文件转换为md格式。
    #  '''

def extract_text_from_pdf(pdf_path):
    text = ""
    # 打开 PDF 文件
    pdf_document = fitz.open(pdf_path)
    # 遍历每一页
    for page_number in range(len(pdf_document)):
        page = pdf_document.load_page(page_number)
        text += page.get_text().replace('\n', ' ')
        # 提取图片并进行 OCR
        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = pdf_document.extract_image(xref)
            image_bytes = base_image["image"]
            image = Image.open(io.BytesIO(image_bytes))
            # 使用 Tesseract OCR 进行文本提取
            text += pytesseract.image_to_string(image,
                                                lang='chi_sim')
#             print(text)
    # 关闭 PDF 文件
    pdf_document.close()
    return text


def extract_text_from_docx(docx_path):
    doc = Document(docx_path)
    text = ''
    for paragraph in doc.paragraphs:
        text += paragraph.text
    return text

def read_wps_content(file_path):
    import win32com.client
    word = win32com.client.Dispatch('Word.Application')
    doc = word.Documents.Open(file_path)
    content = doc.Content.Text
    doc.Close(False)
    word.Quit()
    return content

def extract_text_from_xlsx(xlsx_path):
    wb = xlrd.open_workbook(xlsx_path)
    text = ''
    for sheet in wb.sheets():
        for row in range(sheet.nrows):
            for col in range(sheet.ncols):
                cell_value = str(sheet.cell(row, col).value)
                text += cell_value + '\n'
    return text


def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    index = 1
    text = ''
    # 读取幻灯片的每一页
    for slide in prs.slides:
        # 读取每一板块
        for shape in slide.shapes:
            # print(dir(shape))
            # 是否有文字框
            if shape.has_text_frame:
                # 读文字框的每一段落
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text:
                        # 输出段落文字,也有一些属性,可以用dir查看
                        # print(dir(paragraph))
                        text += paragraph.text
            # 是否有表格
            elif shape.has_table:
                one_table_data = []
                for row in shape.table.rows:  # 读每行
                    row_data = []
                    for cell in row.cells:  # 读一行中的所有单元格
                        c = cell.text
                        row_data.append(c)
                    one_table_data.append(row_data)  # 把每一行存入表
                # 用二维列表输出表格行和列的数据
                print(one_table_data)
                text += one_table_data
            # 是否有图片
            elif isinstance(shape, Picture):
                # shape.image.blob:二进制图像字节流,写入图像文件
                with open(f'{index}.jpg', 'wb') as f:
                    f.write(shape.image.blob)
                    index += 1
                image = Image.open(io.BytesIO(image_bytes))
                # 使用 Tesseract OCR 进行文本提取
                text += pytesseract.image_to_string(image,
                                                lang='chi_sim').replace('', ' ')
    return text

# 定义数据清理函数
def clean_text(text):
    # 删除特殊字符和标点符号
    # text = re.sub(r'[^\w\s]', '', text)
    # 转换为小写
    keywords = ["政府", "环保", "市", "督查", "襄阳", "十堰", "经开区", "招投标"]
    text = replace_keywords(text, keywords, "xx")
    text = text.lower()
    return text


def replace_keywords(text, keywords, repacement):
    for keywod in keywords:
        text = text.replace(keywod, repacement)
    return text

# def traverse_directory(directory):
#     for root, dirs, files in os.walk(directory):
#         for name in files:
#             print(os.path.join(root, name))
#         for name in dirs:
#             print(os.path.join(root, name))
# 提取文本数据并进行清理
def _type2md(path):
    data = []
    # data_directory = './data/1930f329f8ca2000/'  # 数据目录
    data_directory = path
    oa_files = []    
    # 创建个人单个文件的 知识库
    
    for root, dirs, files in os.walk(data_directory):
        for name in files:
            file_path = os.path.join(root, name)
            # print(file_path)
        # PDF 文件
            if name.endswith('.pdf'):
                text = extract_text_from_pdf(file_path)
                relative_path = os.path.relpath(file_path, data_directory)
                output_path = os.path.join(
                    data_directory, 'vertor/'+relative_path.split('\\')[0])
                if not os.path.exists(output_path):
                    os.makedirs(output_path)
                output_path = output_path+'/' + \
                    relative_path.split('\\')[1].split('.')[0] + '.md'
                oa_files.append(output_path)
            # 保存内容到 txt 文件
                with open(output_path, 'w', encoding='utf-8') as output_file:
                    output_file.write(text)
                
        # docx 、 wps文件
            elif name.endswith('.docx') or name.endswith('.wps'):
                
                new_file_path = file_path.replace(
                '\\', '/')
                if name.endswith('.wps'):
                    text = read_wps_content("E:\\vs-python\\AGENTS\\"+new_file_path.split('./')[1].replace('/','\\'))
                else:
                    text = extract_text_from_docx(new_file_path)
                relative_path = os.path.relpath(new_file_path, data_directory)
                output_path = os.path.join(
                    data_directory, 'vertor/' + relative_path.split('\\')[0])
                if not os.path.exists(output_path):
                    os.makedirs(output_path)
            # 保存内容到 txt 文件
                output_path = output_path+'/' + \
                    relative_path.split('\\')[1].split('.')[0] + '.md'
                oa_files.append(output_path)
                with open(output_path, 'w', encoding='utf-8') as output_file:
                    output_file.write(text)
                
        # doc文件  
            elif name.endswith('.doc'):
                import office
                new_file_path = file_path.replace(
                '\\', '/').replace('.doc', '.docx')
                office.word.doc2docx(file_path.replace(
                '\\', '/'), file_path.split('\\')[0])
                text = extract_text_from_docx(new_file_path)
                relative_path = os.path.relpath(new_file_path, data_directory)
                output_path = os.path.join(
                    data_directory, 'vertor/' + relative_path.split('\\')[0])
                if not os.path.exists(output_path):
                    os.makedirs(output_path)
            # 保存内容到 txt 文件
                output_path = output_path+'/' + \
                    relative_path.split('\\')[1].split('.')[0] + '.md'
                oa_files.append(output_path)
                with open(output_path, 'w', encoding='utf-8') as output_file:
                    output_file.write(text)
                
        # xlsx 文件
            elif name.endswith('.xlsx'):
                text = extract_text_from_xlsx(file_path)
                relative_path = os.path.relpath(file_path, data_directory)
                output_path = os.path.join(
                    data_directory, 'vertor/'+relative_path.split('\\')[0])
                if not os.path.exists(output_path):
                    os.makedirs(output_path)
            # 保存内容到 txt 文件
                output_path = output_path+'/' + \
                    relative_path.split('\\')[1].split('.')[0] + '.md'
                oa_files.append(output_path)
                with open(output_path, 'w', encoding='utf-8') as output_file:
                    output_file.write(text)

        # html文件    
            elif name.endswith('.html'):
                import pdfkit
                new_file_path = file_path.replace(
                '\\', '/').replace('.html', '.pdf')
                pdfkit.from_file(file_path.replace('\\', '/'), new_file_path)
                text = extract_text_from_pdf(new_file_path)
                relative_path = os.path.relpath(new_file_path, data_directory)
                output_path = os.path.join(
                    data_directory, 'vertor/'+relative_path.split('\\')[0])
                if not os.path.exists(output_path):
                    os.makedirs(output_path)
                output_path = output_path+'/' + \
                    relative_path.split('\\')[1].split('.')[0] + '.md'
                oa_files.append(output_path)
                # 保存内容到 txt 文件
                with open(output_path, 'w', encoding='utf-8') as output_file:
                    output_file.write(text)
                
            # pptx 文件     
            elif name.endswith('.pptx'):
                text = extract_text_from_pptx(file_path)
                relative_path = os.path.relpath(file_path, data_directory)
                output_path = os.path.join(
                    data_directory, 'vertor/'+relative_path.split('\\')[0])
                if not os.path.exists(output_path):
                    os.makedirs(output_path)
                output_path = output_path+'/' + \
                    relative_path.split('\\')[1].split('.')[0] + '.md'
                oa_files.append(output_path)
                # 保存内容到 txt 文件
                with open(output_path, 'w', encoding='utf-8') as output_file:
                    output_file.write(text)
                
            else:
                continue

            cleaned_text = clean_text(text)
            data.append({'text': cleaned_text, 'file_name': name})

    # if data['msg'] == '文件上传与向量化完成':
    #     print(oa_search_docs(subject,kb_name))
    return oa_files